import streamlit as st
import datetime
import time
import os
import io
import re
import base64
import zipfile
from urllib.parse import urljoin
import xml.etree.ElementTree as ET  # hwpx 처리를 위한 XML 파싱
import pandas as pd
import pythoncom

from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from bs4 import BeautifulSoup, NavigableString

from openai import OpenAI
import pdfplumber  # PDF 처리용

# HWP 처리 관련 (olefile, hwp5html, html2text 등은 사용하지 않음)

client = os.environ.get("OPENAI_API_KEY")

###############################################
# 웹 크롤링 → Markdown 변환 관련 함수
###############################################
def remove_code_block_markers(text: str) -> str:
    text = re.sub(r'^```(?:\w+)?\n', '', text)
    text = re.sub(r'\n```$', '', text)
    return text

def get_rendered_html(url):
    chrome_options = Options()
    chrome_options.add_argument("--headless")
    chrome_options.add_argument("--no-sandbox")
    chrome_options.add_argument("--disable-dev-shm-usage")
    driver = webdriver.Chrome(options=chrome_options)
    driver.get(url)
    time.sleep(3)  # 동적 콘텐츠 로딩 대기
    html = driver.page_source
    driver.quit()
    return html

def extract_element(element, base_url=None):
    if isinstance(element, NavigableString):
        return element.strip()
    if element.name in ['script', 'style']:
        return ""
    if element.name in ['a', 'img']:
        if element.name == 'img' and base_url:
            src = element.get('src')
            if src and not src.startswith('http'):
                element['src'] = urljoin(base_url, src)
        return str(element)
    result = ""
    for child in element.children:
        child_text = extract_element(child, base_url)
        if child_text:
            result += child_text + " "
    return result.strip()

def extract_content(html, target_class=None, base_url=None):
    soup = BeautifulSoup(html, 'html.parser')
    if target_class:
        container = soup.find(class_=target_class)
        element = container if container else soup
    else:
        element = soup
    return extract_element(element, base_url)

def convert_to_markdown(text):
    prompt = (
            "다음 텍스트에서 문맥상 필요없는 부분이나 읽기 어려운 부분, 그리고 javascript 등 코딩 코드를 제거하고 "
            "적절한 헤딩, 단락, 리스트 등을 사용하여 잘 구조화된 Markdown 문서를 만들어줘."
            "단, 마지막에 추가적인 요약이나 정리 멘트는 생략해줘:\n\n" + text
    )
    response = client.chat.completions.create(
        model="gpt-4o-mini",  # 필요에 따라 모델 변경 가능
        messages=[
            {"role": "system", "content": "너는 평문을 Markdown 형식으로 정제하는 전문가야. 응답은 마크다운 형식 그대로여야 해."},
            {"role": "user", "content": prompt}
        ]
    )
    markdown_text = response.choices[0].message.content
    return markdown_text

def auto_download(markdown_text, file_name):
    b64 = base64.b64encode(markdown_text.encode()).decode()
    download_html = f"""
    <html>
      <body>
        <a id="downloadLink" href="data:text/markdown;base64,{b64}" download="{file_name}" style="display:none">Download</a>
        <script>document.getElementById('downloadLink').click();</script>
      </body>
    </html>
    """
    st.components.v1.html(download_html, height=0)

def process_url(url, target_class):
    html = get_rendered_html(url)
    content = extract_content(html, target_class, base_url=url)
    if not content:
        return None, None
    md_text = convert_to_markdown(content)
    md_text = remove_code_block_markers(md_text)

    header_line = None
    for line in md_text.splitlines():
        if line.strip().startswith("#"):
            header_line = line.strip()
            break
    if header_line:
        sanitized_header = re.sub(r'[\\/*?:"<>|]', "", header_line.lstrip("#").strip())
    else:
        sanitized_header = "converted_document"

    now = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
    file_name = f"{sanitized_header}_{now}.md"
    return file_name, md_text

###############################################
# 파일 → Markdown 변환 관련 함수
###############################################
# def extract_text_from_pdf(file):
#     file_bytes = file.read()
#     pdf_io = io.BytesIO(file_bytes)
#     text = ""
#     from pdfminer.layout import LAParams
#     laparams = LAParams(char_margin=2.0, line_margin=0.5, word_margin=0.1, boxes_flow=0.5)
#     with pdfplumber.open(pdf_io, laparams=laparams.__dict__) as pdf:
#         for page in pdf.pages:
#             page_text = page.extract_text()
#             if page_text:
#                 text += page_text + "\n"
#     return text

def fill_missing_cells(table):
    """
    리스트 형태의 테이블(리스트의 리스트)에서, 각 셀의 값이 None 또는 빈 문자열인 경우,
    바로 위 행의 같은 열 값을 채워 넣는 방식으로 보완합니다.
    단, 첫 번째 행(헤더)는 건드리지 않습니다.
    """
    if not table or len(table) < 2:
        return table
    # 첫 번째 행은 헤더로 가정
    for i in range(1, len(table)):
        row = table[i]
        for j in range(len(row)):
            if row[j] in [None, ""]:
                # 바로 위 행에서 값이 있으면 가져오기
                row[j] = table[i - 1][j] if table[i - 1][j] not in [None, ""] else ""
    return table

def extract_text_from_pdf(file):
    import io
    import pdfplumber
    from pdfminer.layout import LAParams
    import pandas as pd

    file_bytes = file.read()
    pdf_io = io.BytesIO(file_bytes)
    text = ""
    laparams = LAParams(char_margin=2.0, line_margin=0.5, word_margin=0.1, boxes_flow=0.5)
    with pdfplumber.open(pdf_io, laparams=laparams.__dict__) as pdf:
        for page in pdf.pages:
            # 추출된 일반 텍스트
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

            # 페이지 내 테이블 추출
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    # 병합된 셀(누락된 값) 보완: 이전 행의 값으로 채워 넣기
                    table = fill_missing_cells(table)
                    # 테이블을 Markdown으로 변환 (첫 행을 헤더로 사용)
                    try:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        markdown_table = df.to_markdown(index=False)
                    except Exception:
                        # DataFrame 변환에 실패하면 단순 문자열 결합 방식 사용
                        markdown_table = "\n".join([" | ".join(map(str, row)) for row in table])
                    text += "\n" + markdown_table + "\n\n"
    return text


def extract_text_from_hwpx(file):
    try:
        file_bytes = file.read()
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            names = z.namelist()
            candidate = next((name for name in names if "section0.xml" in name.lower()), None)
            if not candidate:
                st.error("hwpx 파일에서 section XML을 찾을 수 없습니다.")
                return ""
            xml_content = z.read(candidate)
    except Exception as e:
        st.error(f"hwpx 파일 처리 오류: {e}")
        return ""
    try:
        tree = ET.fromstring(xml_content)
    except Exception as e:
        st.error(f"XML 파싱 오류: {e}")
        return ""
    texts = [elem.text.strip() for elem in tree.iter() if elem.text]
    return " ".join(texts).strip()

def extract_text_from_txt(file):
    try:
        content = file.read()
        if isinstance(content, bytes):
            content = content.decode("utf-8")
    except Exception as e:
        st.error(f"TXT 파일 추출 오류: {e}")
        content = ""
    return content

def extract_text_from_pptx(file):
    try:
        from pptx import Presentation
    except ImportError:
        st.error("pptx 파일 처리를 위해 python-pptx 라이브러리가 필요합니다. pip install python-pptx")
        return ""
    try:
        file.seek(0)
        ppt = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"PPTX 파일 추출 오류: {e}")
        return ""

def extract_text_from_xlsx(file):
    try:
        file_bytes = file.read()
        from openpyxl import load_workbook
        wb = load_workbook(filename=io.BytesIO(file_bytes), data_only=True)
        result_text = ""
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            # 병합된 셀 처리: 각 병합 영역에 대해 unmerge 후, 모든 셀에 왼쪽 상단 값 채우기
            merged_ranges = list(ws.merged_cells.ranges)
            for merged_range in merged_ranges:
                top_left_value = ws.cell(merged_range.min_row, merged_range.min_col).value
                # 병합 영역 해제
                ws.unmerge_cells(range_string=str(merged_range))
                # 해제된 각 셀에 값을 채워 넣기
                for row in range(merged_range.min_row, merged_range.max_row + 1):
                    for col in range(merged_range.min_col, merged_range.max_col + 1):
                        ws.cell(row=row, column=col).value = top_left_value
            # 시트의 데이터를 리스트 형태로 추출 (첫 행은 컬럼명)
            data = ws.values
            columns = next(data)
            df = pd.DataFrame(data, columns=columns)
            try:
                markdown_table = df.to_markdown(index=False)
            except Exception:
                markdown_table = df.to_string(index=False)
            result_text += f"Sheet: {sheet_name}\n{markdown_table}\n\n"
        return result_text
    except Exception as e:
        st.error(f"XLSX 파일 추출 오류: {e}")
        return ""


def extract_text_from_file(file):
    ext = os.path.splitext(file.name)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file)
    elif ext == ".hwp":
        st.error("HWP 파일 추출 기능이 구현되어 있지 않습니다.")
        return ""
    elif ext == ".hwpx":
        return extract_text_from_hwpx(file)
    elif ext == ".txt":
        return extract_text_from_txt(file)
    elif ext == ".pptx":
        return extract_text_from_pptx(file)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file)
    else:
        st.error("지원되지 않는 파일 형식입니다. (PDF, HWP, HWPX, TXT, PPTX, XLSX 파일만 지원)")
        return ""

def split_text_with_overlap(text, chunk_size, overlap):
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunks.append(" ".join(words[start:end]))
        start = end - overlap
    return chunks

def remove_code_fence(markdown_text):
    markdown_text = markdown_text.strip()
    if markdown_text.startswith("```"):
        lines = markdown_text.splitlines()
        if lines[0].startswith("```"):
            if lines[-1].strip().startswith("```"):
                markdown_text = "\n".join(lines[1:-1]).strip()
            else:
                markdown_text = "\n".join(lines[1:]).strip()
    return markdown_text

def remove_isolated_code_fences(text):
    return re.sub(r"(?m)^\s*```(markdown)?\s*$", "", text)

def convert_chunk_to_markdown(current_chunk, previous_md=None):
    context = ""
    if previous_md:
        prev_words = previous_md.split()
        context = " ".join(prev_words[-100:]) + "\n\n"
    prompt = (
        "**현재청크**를 잘 구조화된 Markdown 문서로 변환해줘. "
        "Markdown 문서의 문맥을 분석하여 어색한 부분들은 모두 삭제해줘. "
        "**이전 마크다운 청크**는 포함하지 않고 **현재청크**만 정제해서 만들어줘. "
        f"**이전 마크다운 청크**: {context}\n\n"
        f"**현재청크**: {current_chunk}"
    )
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "너는 텍스트 변환 전문가야. 주어진 청크의 모든 세부 정보를 보존하면서 Markdown 문서로 변환해."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.3
    )
    markdown_text = response.choices[0].message.content
    return remove_code_fence(markdown_text)

###############################################
# Markdown → TXT 변환 관련 함수
###############################################
def transform_level1_header(md_content):
    return re.sub(r'^( *)#(?!#)\s+', r'\1@!@ ', md_content, flags=re.MULTILINE)

###############################################
# 통합 Streamlit 앱 메인 함수
###############################################
def main():
    st.title("통합 변환기")
    st.write("웹 크롤링, 파일 변환 및 Markdown → TXT 변환 기능을 제공합니다.")

    # 탭 순서: [Web Crawl, File to Markdown, HWP to HWPX 변환, Markdown to TXT 변환]
    tabs = st.tabs([
        "Web Crawl to Markdown",
        "File to Markdown",
        "Markdown to TXT 변환"
    ])

    # 1. Web Crawl to Markdown 탭
    with tabs[0]:
        st.header("웹 크롤링 → Markdown 변환")
        st.write(
            "여러 URL을 입력하고 (각 줄에 하나씩) 옵션으로 특정 클래스명을 지정하면, "
            "각 URL의 동적 페이지를 렌더링하여 텍스트와 <a>, <img> 등의 HTML 마크업을 그대로 추출한 후, "
            "OpenAI를 통해 Markdown으로 정제한 결과를 다운로드합니다."
        )
        urls_text = st.text_area("크롤링할 URL들을 입력하세요 (각 줄에 하나씩):", "https://example.com")
        target_class = st.text_input("특정 클래스명을 입력하세요 (옵션):", "")
        if st.button("문서 생성 (Web Crawl)"):
            urls = [line.strip() for line in urls_text.splitlines() if line.strip()]
            if not urls:
                st.error("적어도 하나의 URL을 입력하세요.")
            else:
                for url in urls:
                    st.info(f"처리 중: {url}")
                    try:
                        file_name, md_text = process_url(url, target_class.strip() if target_class.strip() else None)
                        if md_text is None:
                            st.warning(f"{url}에서 콘텐츠를 추출하지 못했습니다.")
                            continue
                        st.success(f"{url} -> 생성된 파일: {file_name}")
                        auto_download(md_text, file_name)
                        st.markdown("---")
                        time.sleep(1)
                    except Exception as e:
                        st.error(f"{url} 처리 중 오류 발생: {e}")

    # 2. File to Markdown 탭
    with tabs[1]:
        st.header("파일 → Markdown 변환")
        st.write(
            "PDF, HWP, HWPX, TXT, PPTX, XLSX 파일을 업로드하면 전체 텍스트를 추출하고 청크 단위로 분할한 후, "
            "각 청크를 OpenAI를 통해 Markdown 형식으로 정제합니다. "
            "파일별로 Markdown 및 TXT 다운로드 버튼과 전체 결과 ZIP 파일 다운로드 기능을 제공합니다.\n\n"
            "※ TXT 파일만 업로드한 경우 단순 합본 TXT 파일을 생성합니다."
        )
        overall_download_container = st.container()
        uploaded_files = st.file_uploader(
            "파일 업로드 (PDF, HWP, HWPX, TXT, PPTX, XLSX)",
            type=["pdf", "hwp", "hwpx", "txt", "pptx", "xlsx"],
            accept_multiple_files=True
        )
        if uploaded_files:
            all_txt = all(os.path.splitext(f.name)[1].lower() == ".txt" for f in uploaded_files)
            if all_txt:
                merged_text = ""
                for uploaded_file in uploaded_files:
                    st.write(f"### {uploaded_file.name} 파일 처리 중...")
                    with st.spinner("TXT 파일 내용 읽는 중..."):
                        content = extract_text_from_file(uploaded_file)
                    merged_text += content
                st.success("TXT 파일 합본 생성 완료!")
                st.download_button(
                    label="합본 TXT 파일 다운로드",
                    data=merged_text,
                    file_name="merged_text.txt",
                    mime="text/plain",
                    key="merged_txt_only"
                )
            else:
                results_md = []
                results_txt = []
                for uploaded_file in uploaded_files:
                    st.write(f"### 파일 처리 중: {uploaded_file.name}")
                    with st.spinner("파일에서 텍스트 추출 중..."):
                        full_text = extract_text_from_file(uploaded_file)
                    if not full_text:
                        st.error(f"{uploaded_file.name} 파일의 텍스트 추출에 실패했습니다.")
                        continue
                    chunks = split_text_with_overlap(full_text, chunk_size=1000, overlap=0)
                    st.info(f"{uploaded_file.name} - 총 청크 개수: {len(chunks)}")
                    refined_markdowns = []
                    previous_md = None
                    for idx, chunk in enumerate(chunks):
                        with st.spinner(f"{uploaded_file.name} - 청크 {idx + 1}/{len(chunks)} 처리 중..."):
                            md_chunk = convert_chunk_to_markdown(chunk, previous_md)
                            refined_markdowns.append(md_chunk)
                            previous_md = md_chunk
                        time.sleep(1)
                    final_markdown = "\n\n---\n\n".join(refined_markdowns)
                    final_markdown = remove_isolated_code_fences(final_markdown)
                    base_name = os.path.splitext(uploaded_file.name)[0]
                    md_file_name = f"{base_name}.md"
                    results_md.append((md_file_name, final_markdown))
                    st.success(f"{uploaded_file.name} Markdown 문서 생성 완료: {md_file_name}")
                    st.code(final_markdown, language="markdown")
                    st.download_button(
                        label=f"{md_file_name} 다운로드",
                        data=final_markdown,
                        file_name=md_file_name,
                        mime="text/markdown",
                        key=f"md_{base_name}"
                    )
                    transformed_txt = transform_level1_header(final_markdown)
                    txt_file_name = f"{base_name}.txt"
                    results_txt.append((txt_file_name, transformed_txt))
                    st.download_button(
                        label=f"{txt_file_name} 다운로드 (TXT)",
                        data=transformed_txt,
                        file_name=txt_file_name,
                        mime="text/plain",
                        key=f"txt_{base_name}"
                    )
                    st.markdown("---")
                if results_md:
                    zip_buffer_md = io.BytesIO()
                    with zipfile.ZipFile(zip_buffer_md, "w") as zf:
                        for fname, md_content in results_md:
                            zf.writestr(fname, md_content)
                    zip_buffer_md.seek(0)
                else:
                    zip_buffer_md = None
                if results_txt:
                    merged_txt = "".join(
                        [f"\n{txt_content}" for fname, txt_content in results_txt]
                    )
                    zip_buffer_txt = io.BytesIO()
                    with zipfile.ZipFile(zip_buffer_txt, "w") as zf:
                        for fname, txt_content in results_txt:
                            zf.writestr(fname, txt_content)
                    zip_buffer_txt.seek(0)
                else:
                    zip_buffer_txt = None
                    merged_txt = ""
                with overall_download_container:
                    cols = st.columns([4, 2])
                    with cols[1]:
                        if zip_buffer_md:
                            st.download_button(
                                label="모든 Markdown 파일 ZIP 다운로드",
                                data=zip_buffer_md,
                                file_name="converted_markdown_files.zip",
                                mime="application/zip",
                                key="overall_zip_md"
                            )
                        if zip_buffer_txt:
                            st.download_button(
                                label="모든 TXT 파일 ZIP 다운로드",
                                data=zip_buffer_txt,
                                file_name="converted_txt_files.zip",
                                mime="application/zip",
                                key="overall_zip_txt"
                            )
                        if merged_txt:
                            st.download_button(
                                label="모든 TXT 병합 파일 다운로드",
                                data=merged_txt,
                                file_name="merged_txt.txt",
                                mime="text/plain",
                                key="merged_txt"
                            )


    # 4. Markdown → TXT 변환 탭
    with tabs[2]:
        st.header("Markdown → TXT 변환")
        st.write("업로드한 Markdown(.md) 파일에서 레벨 1 헤더(#)를 '@!@'로 치환하여 텍스트(.txt) 파일로 변환합니다.")
        uploaded_md_files = st.file_uploader("Markdown 파일 업로드", type=["md"], accept_multiple_files=True)
        if uploaded_md_files:
            results = []
            for uploaded_file in uploaded_md_files:
                st.write(f"### 처리 중: {uploaded_file.name}")
                try:
                    md_content = uploaded_file.read().decode("utf-8")
                except Exception as e:
                    st.error(f"{uploaded_file.name} 파일 읽기 오류: {e}")
                    continue
                transformed_content = transform_level1_header(md_content)
                base_name = os.path.splitext(uploaded_file.name)[0]
                file_name = f"{base_name}.txt"
                results.append((file_name, transformed_content))
                st.success(f"{uploaded_file.name} → {file_name} 변환 완료")
                st.download_button(
                    label=f"{file_name} 다운로드",
                    data=transformed_content,
                    file_name=file_name,
                    mime="text/plain"
                )
                st.code(transformed_content, language="plaintext")
                st.markdown("---")
            if len(results) > 1:
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, "w") as zip_file:
                    for fname, content in results:
                        zip_file.writestr(fname, content)
                zip_buffer.seek(0)
                st.download_button(
                    label="모든 TXT 파일 ZIP 다운로드",
                    data=zip_buffer,
                    file_name="converted_txt_files.zip",
                    mime="application/zip"
                )

if __name__ == "__main__":
    main()
