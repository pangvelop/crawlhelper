import os
import io
import zipfile
import streamlit as st
import pandas as pd
import xml.etree.ElementTree as ET
import pdfplumber
from pdfminer.layout import LAParams

def handle_file_to_md_tab():
    st.header("íŒŒì¼ â†’ Markdown ë³€í™˜ê¸°")
    uploaded_file = st.file_uploader("íŒŒì¼ ì—…ë¡œë“œ (PDF, HWPX, PPTX, TXT, XLSX)", type=["pdf", "hwpx", "pptx", "txt", "xlsx"])

    if uploaded_file is not None:
        with st.spinner("íŒŒì¼ ì²˜ë¦¬ ì¤‘..."):
            extracted_text = extract_text_from_file(uploaded_file)

        if extracted_text.strip():
            st.subheader("ğŸ“„ ì¶”ì¶œëœ Markdown")
            st.text_area("Markdown ê²°ê³¼", extracted_text, height=400)

            st.download_button("ğŸ’¾ Markdown ë‹¤ìš´ë¡œë“œ", data=extracted_text, file_name="converted.md")
        else:
            st.warning("ë³€í™˜ ê²°ê³¼ê°€ ë¹„ì–´ ìˆìŠµë‹ˆë‹¤.")


def extract_text_from_pdf(file):
    file_bytes = file.read()
    pdf_io = io.BytesIO(file_bytes)
    text = ""
    laparams = LAParams(char_margin=2.0, line_margin=0.5, word_margin=0.1, boxes_flow=0.5)
    with pdfplumber.open(pdf_io, laparams=laparams.__dict__) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

            tables = page.extract_tables()
            if tables:
                for table in tables:
                    table = fill_missing_cells(table)
                    try:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        markdown_table = df.to_markdown(index=False)
                    except Exception:
                        markdown_table = "\n".join([" | ".join(map(str, row)) for row in table])
                    text += "\n" + markdown_table + "\n\n"
    return text


def fill_missing_cells(table):
    if not table or len(table) < 2:
        return table
    for i in range(1, len(table)):
        row = table[i]
        for j in range(len(row)):
            if row[j] in [None, ""]:
                row[j] = table[i - 1][j] if table[i - 1][j] not in [None, ""] else ""
    return table


def extract_text_from_hwpx(file):
    try:
        file_bytes = file.read()
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            names = z.namelist()
            candidate = next((name for name in names if "section0.xml" in name.lower()), None)
            if not candidate:
                st.error("hwpx íŒŒì¼ì—ì„œ section XMLì„ ì°¾ì„ ìˆ˜ ì—†ìŠµë‹ˆë‹¤.")
                return ""
            xml_content = z.read(candidate)
    except Exception as e:
        st.error(f"hwpx íŒŒì¼ ì²˜ë¦¬ ì˜¤ë¥˜: {e}")
        return ""
    try:
        tree = ET.fromstring(xml_content)
    except Exception as e:
        st.error(f"XML íŒŒì‹± ì˜¤ë¥˜: {e}")
        return ""
    texts = [elem.text.strip() for elem in tree.iter() if elem.text]
    return " ".join(texts).strip()


def extract_text_from_txt(file):
    try:
        content = file.read()
        if isinstance(content, bytes):
            content = content.decode("utf-8")
        return content
    except Exception as e:
        st.error(f"TXT íŒŒì¼ ì¶”ì¶œ ì˜¤ë¥˜: {e}")
        return ""


def extract_text_from_pptx(file):
    try:
        from pptx import Presentation
    except ImportError:
        st.error("python-pptxê°€ ì„¤ì¹˜ë˜ì–´ ìˆì§€ ì•ŠìŠµë‹ˆë‹¤. pip install python-pptx")
        return ""
    try:
        file.seek(0)
        ppt = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"PPTX íŒŒì¼ ì¶”ì¶œ ì˜¤ë¥˜: {e}")
        return ""


def extract_text_from_xlsx(file):
    try:
        from openpyxl import load_workbook
        file_bytes = file.read()
        wb = load_workbook(filename=io.BytesIO(file_bytes), data_only=True)
        result_text = ""
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            merged_ranges = list(ws.merged_cells.ranges)
            for merged_range in merged_ranges:
                top_left_value = ws.cell(merged_range.min_row, merged_range.min_col).value
                ws.unmerge_cells(range_string=str(merged_range))
                for row in range(merged_range.min_row, merged_range.max_row + 1):
                    for col in range(merged_range.min_col, merged_range.max_col + 1):
                        ws.cell(row=row, column=col).value = top_left_value
            data = ws.values
            columns = next(data)
            df = pd.DataFrame(data, columns=columns)
            try:
                markdown_table = df.to_markdown(index=False)
            except Exception:
                markdown_table = df.to_string(index=False)
            result_text += f"Sheet: {sheet_name}\n{markdown_table}\n\n"
        return result_text
    except Exception as e:
        st.error(f"XLSX íŒŒì¼ ì¶”ì¶œ ì˜¤ë¥˜: {e}")
        return ""


def extract_text_from_file(file):
    ext = os.path.splitext(file.name)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file)
    elif ext == ".hwp":
        st.error("HWP íŒŒì¼ ì¶”ì¶œ ê¸°ëŠ¥ì€ ì§€ì›ë˜ì§€ ì•ŠìŠµë‹ˆë‹¤.")
        return ""
    elif ext == ".hwpx":
        return extract_text_from_hwpx(file)
    elif ext == ".txt":
        return extract_text_from_txt(file)
    elif ext == ".pptx":
        return extract_text_from_pptx(file)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file)
    else:
        st.error("ì§€ì›ë˜ì§€ ì•ŠëŠ” íŒŒì¼ í˜•ì‹ì…ë‹ˆë‹¤. (PDF, HWPX, TXT, PPTX, XLSX)")
        return ""
